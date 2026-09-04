"""Multi-format document loader.

Accepts markdown, plaintext, HTML and PDF files and normalizes each into a
`LoadedDocument`: clean plaintext plus metadata (source file, section
heading structure, page number for PDFs).

Design note on "section heading" across formats
-------------------------------------------------
Markdown already uses '#'..'######' as its heading convention. To let every
downstream chunker (in particular the structure-aware chunker) treat all
four formats uniformly, HTML headings (<h1>..<h6>) are converted to the same
'#'..'######' markdown-style prefix during normalization. Plaintext has no
heading convention and is passed through unchanged (one untitled section).

PDF is intentionally NOT given synthetic headings. `pypdf`'s text extraction
does not expose font size/weight, so any "is this line a heading" heuristic
built on text alone is unreliable (false positives on short lines, ALL-CAPS
labels, etc.) and would silently corrupt structure-aware chunking. Instead,
PDFs keep their one genuinely reliable structural signal: page number,
tracked per page via `LoadedDocument.pages`.

Word documents (.docx) and Presentations (.pptx) carry explicit structure,
so they map directly to markdown headings (`#` through `######`). .pptx
slides also track `page_number` just like PDFs.
"""
from __future__ import annotations

import re
from pathlib import Path

from bs4 import BeautifulSoup
from pypdf import PdfReader

from rag_api.domain.models import LoadedDocument, PageText

SUPPORTED_EXTENSIONS = {".md", ".markdown", ".txt", ".html", ".htm", ".pdf", ".docx", ".pptx", ".xlsx"}


class UnsupportedFileTypeError(ValueError):
    pass


class EmptyDocumentError(ValueError):
    pass


def load_document(path: str | Path, llm_client=None, image_store=None) -> LoadedDocument:
    """Load and normalize a single file. Raises UnsupportedFileTypeError /
    EmptyDocumentError / FileNotFoundError as appropriate."""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"No such file: {path}")
        
    from rag_api.core.settings import get_settings
    settings = get_settings()

    suffix = path.suffix.lower()
    if suffix in (".md", ".markdown"):
        doc = _load_markdown(path)
    elif suffix == ".txt":
        doc = _load_text(path)
    elif suffix in (".html", ".htm"):
        doc = _load_html(path, settings, llm_client, image_store)
    elif suffix == ".pdf":
        doc = _load_pdf(path, llm_client, image_store)
    elif suffix == ".docx":
        doc = _load_docx(path, settings, llm_client, image_store)
    elif suffix == ".pptx":
        doc = _load_pptx(path)
    elif suffix == ".xlsx":
        doc = _load_xlsx(path)
    else:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '{suffix}' for {path.name}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )

    if not doc.text.strip():
        raise EmptyDocumentError(f"{path.name} contains no extractable text")
    return doc


def load_documents(paths: list[str | Path], llm_client=None, image_store=None) -> tuple[list[LoadedDocument], list[dict]]:
    """Load many files. Never raises for a single bad file — collects errors
    instead so one broken upload doesn't abort a whole ingestion batch.
    Returns (loaded_documents, errors) where each error is
    {"source_file": ..., "error": ...}."""
    loaded: list[LoadedDocument] = []
    errors: list[dict] = []
    for p in paths:
        try:
            loaded.append(load_document(p, llm_client=llm_client, image_store=image_store))
        except Exception as exc:  # noqa: BLE001 - intentionally broad, collected not raised
            errors.append({"source_file": str(Path(p).name), "error": str(exc)})
    return loaded, errors


def _load_markdown(path: Path) -> LoadedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    return LoadedDocument(source_file=path.name, format="md", text=text)


def _load_text(path: Path) -> LoadedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    from rag_api.core.settings import get_settings
    settings = get_settings()
    if settings.text_heading_detection_enabled:
        lines = []
        for line in text.splitlines():
            if not line.strip():
                lines.append(line)
                continue
            # Heuristics: ALL CAPS (and > 3 chars), or starts with number e.g. "1. Introduction"
            import re
            is_all_caps = line.isupper() and len(line.strip()) > 3 and len(line.split()) < 10
            is_numbered = bool(re.match(r"^\d+(\.\d+)*\.?\s+[A-Z]", line.strip())) and len(line.split()) < 15
            if is_all_caps or is_numbered:
                lines.append(f"## {line}")
            else:
                lines.append(line)
        text = "\n".join(lines)
    return LoadedDocument(source_file=path.name, format="txt", text=text)


_HEADING_TAGS = {"h1": "#", "h2": "##", "h3": "###", "h4": "####", "h5": "#####", "h6": "######"}
_BLOCK_TAGS = _HEADING_TAGS.keys() | {"p", "li", "blockquote", "td", "th", "pre", "caption"}


def _load_html(path: Path, settings, llm_client, image_store) -> LoadedDocument:
    import base64
    import urllib.request
    from rag_api.domain.models import ExtractedImage
    
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")

    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    lines: list[str] = []
    extracted_images = []
    
    # To keep reading order, we iterate over all tags we care about
    for el in soup.find_all(list(_BLOCK_TAGS) + ["img"]):
        if el.name == "img":
            src = el.get("src", "")
            img_bytes = None
            content_type = "image/png"
            if src.startswith("data:image/"):
                try:
                    header, encoded = src.split(",", 1)
                    img_bytes = base64.b64decode(encoded)
                    content_type = header.split(";")[0][5:]
                except Exception:
                    pass
            elif settings.fetch_remote_html_images and src.startswith("http"):
                try:
                    import httpx
                    resp = httpx.get(src, timeout=5.0)
                    if resp.status_code == 200:
                        img_bytes = resp.content
                        content_type = resp.headers.get("content-type", "image/png")
                except Exception:
                    pass
                    
            if img_bytes:
                try:
                    derived_text, extracted_type = _extract_image_text(img_bytes, settings, llm_client)
                    if extracted_type == "image_untranscribed" and not settings.image_captioning_enabled:
                        continue
                        
                    image_hash = image_store.put(img_bytes, content_type) if image_store else "no_store"
                    extracted_images.append(ExtractedImage(
                        image_hash=image_hash,
                        page_number=None,
                        bbox=None,
                        reading_order_index=len(extracted_images),
                        content_type=extracted_type,
                        derived_text=derived_text
                    ))
                    lines.append(f"\n\n<!--IMG:{image_hash}-->\n\n{derived_text}\n\n<!--/IMG-->\n\n")
                except Exception:
                    pass
        else:
            content = " ".join(el.get_text(separator=" ", strip=True).split())
            if not content:
                continue
            prefix = _HEADING_TAGS.get(el.name)
            lines.append(f"{prefix} {content}" if prefix else content)

    text = "\n\n".join(lines)
    return LoadedDocument(source_file=path.name, format="html", text=text, images=extracted_images)


from rag_api.core.settings import get_settings
from collections import defaultdict

def _load_pdf(path: Path, llm_client=None, image_store=None) -> LoadedDocument:
    settings = get_settings()
    if settings.image_indexing_enabled:
        return _load_pdf_pypdfium2(path, settings, llm_client, image_store)
    if settings.pdf_extraction_backend == "pdfplumber" or settings.pdf_extraction_backend == "pymupdf":
        return _load_pdf_pymupdf(path, settings)
    elif settings.pdf_extraction_backend == "pymupdf":
        raise NotImplementedError("PyMuPDF extraction not yet implemented.")
    else:
        # Fallback to old pypdf logic
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages: list[PageText] = []
        for i, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                pages.append(PageText(page_number=i, text=page_text))
        full_text = "\n\n".join(p.text for p in pages)
        return LoadedDocument(source_file=path.name, format="pdf", text=full_text, pages=pages)

def _load_pdf_pdfplumber(path: Path, settings) -> LoadedDocument:
    import pdfplumber
    pages: list[PageText] = []
    
    with pdfplumber.open(path) as pdf:
        # Pass 1: compute doc-wide median font size to anchor the heading heuristic
        all_sizes = []
        for page in pdf.pages:
            for char in page.chars:
                if char.get("size"):
                    all_sizes.append(char["size"])
        
        if not all_sizes:
            median_size = 12.0
        else:
            all_sizes.sort()
            median_size = all_sizes[len(all_sizes) // 2]
            
        for i, page in enumerate(pdf.pages, start=1):
            # Extract tables if enabled
            tables_data = []
            if settings.pdf_table_extraction_enabled:
                tables = page.find_tables()
                for table in tables:
                    table_md = _format_table(table.extract())
                    bbox = table.bbox
                    tables_data.append({"bbox": bbox, "md": table_md})
            
            # Extract text words
            words = page.extract_words(extra_attrs=["size", "fontname"])
            
            # Group words by "top" position to form lines
            # A simple bucket by rounding to nearest integer
            lines = defaultdict(list)
            for word in words:
                top = round(word["top"] / 2) * 2  # slight leniency in vertical alignment
                lines[top].append(word)
                
            sorted_tops = sorted(lines.keys())
            
            elements = []
            for top in sorted_tops:
                line_words = lines[top]
                line_words.sort(key=lambda w: w["x0"])
                
                avg_size = sum(w["size"] for w in line_words) / len(line_words)
                is_bold = any("Bold" in str(w.get("fontname", "")) or "Black" in str(w.get("fontname", "")) for w in line_words)
                text = " ".join(w["text"] for w in line_words)
                
                # Skip if line is inside a table
                in_table = False
                for t in tables_data:
                    t_bbox = t["bbox"]
                    if t_bbox[1] - 2 <= top <= t_bbox[3] + 2:
                        in_table = True
                        break
                
                if in_table:
                    continue
                
                prefix = ""
                ratio = avg_size / median_size if median_size else 1.0
                if ratio > settings.pdf_heading_font_ratio or is_bold:
                    if ratio > 1.5: level = 1
                    elif ratio > 1.25: level = 2
                    else: level = 3
                    level = min(level, settings.pdf_max_heading_levels)
                    prefix = "#" * level + " "
                
                elements.append((top, "text", prefix + text))
                
            for t in tables_data:
                elements.append((t["bbox"][1], "table", t["md"]))
                
            elements.sort(key=lambda x: x[0])
            
            page_text_content = "\n\n".join(e[2] for e in elements).strip()
            if page_text_content:
                pages.append(PageText(page_number=i, text=page_text_content))
                
    full_text = "\n\n".join(p.text for p in pages)
    return LoadedDocument(source_file=path.name, format="pdf", text=full_text, pages=pages)

def _format_table(table_data) -> str:
    if not table_data:
        return ""
    rows = [r for r in table_data if any(c for c in r)]
    if not rows:
        return ""
    lines = []
    for i, row in enumerate(rows):
        cleaned_row = [" ".join(str(cell).split()) if cell else "" for cell in row]
        lines.append(" | ".join(cleaned_row))
        if i == 0:
            lines.append(" | ".join("---" for _ in row))
    return "\n".join(lines)

def _docx_table_to_markdown(table) -> str:
    lines = []
    for i, row in enumerate(table.rows):
        row_text = " | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells)
        lines.append(row_text)
        if i == 0:
            lines.append(" | ".join("---" for _ in row.cells))
    return "\n".join(lines)


def _load_docx(path: Path, settings, llm_client, image_store) -> LoadedDocument:
    import docx
    from rag_api.domain.models import ExtractedImage
    
    document = docx.Document(str(path))
    lines = []
    extracted_images = []
    
    # Simple extraction of images inside runs
    for para in document.paragraphs:
        for run in para.runs:
            # We don't have direct python-docx API for inline shapes' bytes easily without walking drawing XML
            # Let's use a simpler heuristic for MVP if python-docx provides images, else just skip
            # python-docx has document.inline_shapes for all shapes, but no easy mapping to paragraphs
            pass
            
        text = para.text.strip()
        if text:
            style = para.style.name if para.style else ""
            m = re.match(r"Heading (\d)", style)
            if m:
                level = min(int(m.group(1)), 6)
                lines.append(f"{'#' * level} {text}")
            else:
                lines.append(text)
                
    # Document-level inline shapes (out of order, but it works for MVP)
    for shape in document.inline_shapes:
        try:
            blip = shape._inline.graphic.graphicData.pic.blipFill.blip
            rId = blip.embed
            image_part = document.part.related_parts[rId]
            img_bytes = image_part.blob
            
            derived_text, content_type = _extract_image_text(img_bytes, settings, llm_client)
            if content_type == "image_untranscribed" and not settings.image_captioning_enabled:
                continue
                
            image_hash = image_store.put(img_bytes, image_part.content_type) if image_store else "no_store"
            extracted_images.append(ExtractedImage(
                image_hash=image_hash,
                page_number=None,
                bbox=None,
                reading_order_index=len(extracted_images),
                content_type=content_type,
                derived_text=derived_text
            ))
            lines.append(f"\n\n<!--IMG:{image_hash}-->\n\n{derived_text}\n\n<!--/IMG-->\n\n")
        except Exception:
            pass

    for table in document.tables:
        lines.append(_docx_table_to_markdown(table))
        
    return LoadedDocument(source_file=path.name, format="docx", text="\n\n".join(lines), images=extracted_images)


def _load_pptx(path: Path) -> LoadedDocument:
    from pptx import Presentation
    prs = Presentation(str(path))
    pages: list[PageText] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else None
        body_lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape != slide.shapes.title and shape.text_frame.text.strip()
        ]
        text = ("# " + title + "\n\n" if title else "") + "\n\n".join(body_lines)
        if text.strip():
            pages.append(PageText(page_number=i, text=text))
    full_text = "\n\n".join(p.text for p in pages)
    return LoadedDocument(source_file=path.name, format="pptx", text=full_text, pages=pages)


def _load_xlsx(path: Path) -> LoadedDocument:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    sections = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        header, *body = rows
        header_line = " | ".join(str(c) if c is not None else "" for c in header)
        sep_line = " | ".join("---" for _ in header)
        body_lines = [" | ".join(str(c) if c is not None else "" for c in r) for r in body]
        table_md = "\n".join([header_line, sep_line, *body_lines])
        sections.append(f"## {sheet.title}\n\n{table_md}")
    return LoadedDocument(source_file=path.name, format="xlsx", text="\n\n".join(sections))


def _extract_image_text(image_bytes: bytes, settings, llm_client) -> tuple[str, str]:
    import pytesseract
    from PIL import Image
    from io import BytesIO
    try:
        ocr_text = pytesseract.image_to_string(Image.open(BytesIO(image_bytes))).strip()
        if len(ocr_text.split()) >= settings.min_ocr_words_before_caption_fallback:
            return ocr_text, "image_ocr"
    except Exception:
        ocr_text = ""
        
    if settings.image_captioning_enabled and llm_client:
        try:
            caption = llm_client.describe_image(image_bytes, "image/png", "Describe this image factually.")
            return caption, "image_caption"
        except Exception:
            pass
            
    return ocr_text, "image_untranscribed"

def _load_pdf_pypdfium2(path: Path, settings, llm_client, image_store) -> LoadedDocument:
    import pypdfium2 as pdfium
    import pytesseract
    from PIL import Image
    from io import BytesIO
    from rag_api.domain.models import PageText, ExtractedImage
    
    pages: list[PageText] = []
    extracted_images = []
    pdf = pdfium.PdfDocument(str(path))
    
    for i in range(len(pdf)):
        page = pdf[i]
        page_number = i + 1
        text_page = page.get_textpage()
        native_text = text_page.get_text_bounded().strip()
        
        extraction_method = "native"
        final_text = native_text
        
        if len(native_text) < settings.scanned_page_text_threshold:
            bitmap = page.render(scale=settings.ocr_dpi/72)
            pil_image = bitmap.to_pil()
            ocr_text = pytesseract.image_to_string(pil_image).strip()
            if ocr_text:
                final_text = ocr_text
                extraction_method = "ocr"
        
        images_on_page = []
        for obj in page.get_objects():
            if isinstance(obj, pdfium.PdfImage):
                try:
                    bitmap = obj.get_bitmap()
                    if not bitmap: continue
                    pil_image = bitmap.to_pil()
                    img_bytes_io = BytesIO()
                    pil_image.save(img_bytes_io, format="PNG")
                    img_bytes = img_bytes_io.getvalue()
                    
                    derived_text, content_type = _extract_image_text(img_bytes, settings, llm_client)
                    if content_type == "image_untranscribed" and not settings.image_captioning_enabled:
                        continue # Skip entirely if we can't extract it and captioning is off
                        
                    image_hash = image_store.put(img_bytes, "image/png") if image_store else "no_store"
                    
                    extracted_images.append(ExtractedImage(
                        image_hash=image_hash,
                        page_number=page_number,
                        bbox=None, # pypdfium2 bbox math skipped for MVP
                        reading_order_index=len(images_on_page),
                        content_type=content_type,
                        derived_text=derived_text,
                        ocr_confidence=None
                    ))
                    
                    images_on_page.append(f"\n\n<!--IMG:{image_hash}-->\n\n{derived_text}\n\n<!--/IMG-->\n\n")
                except Exception as e:
                    pass
                    
        if images_on_page:
            final_text += "".join(images_on_page)
            
        if final_text.strip():
            pages.append(PageText(page_number=page_number, text=final_text.strip(), extraction_method=extraction_method))
            
    full_text = "\n\n".join(p.text for p in pages)
    return LoadedDocument(source_file=path.name, format="pdf", text=full_text, pages=pages, images=extracted_images)


def _load_pdf_pymupdf(path: Path, settings) -> LoadedDocument:
    import fitz  # PyMuPDF
    from rag_api.domain.models import PageText
    
    pages = []
    doc = fitz.open(str(path))
    
    for i in range(len(doc)):
        page = doc[i]
        
        # 1) Get text using PyMuPDF dict extraction to grab fonts
        text_content = ""
        blocks = page.get_text("dict").get("blocks", [])
        
        # Pass 1: compute doc-wide median font size for this page
        all_sizes = []
        for b in blocks:
            if b.get("type") == 0:  # text block
                for l in b.get("lines", []):
                    for s in l.get("spans", []):
                        if s.get("size"):
                            all_sizes.append(s["size"])
        
        median_size = 12.0
        if all_sizes:
            all_sizes.sort()
            median_size = all_sizes[len(all_sizes) // 2]
            
        elements = []
        for b in blocks:
            if b.get("type") == 0:
                for l in b.get("lines", []):
                    for s in l.get("spans", []):
                        text = s.get("text", "").strip()
                        if not text: continue
                        size = s.get("size", 12.0)
                        flags = s.get("flags", 0)
                        
                        is_bold = bool(flags & 2 ** 4)
                        ratio = size / median_size if median_size else 1.0
                        
                        prefix = ""
                        if ratio > settings.pdf_heading_font_ratio or is_bold:
                            if ratio > 1.5: level = 1
                            elif ratio > 1.25: level = 2
                            else: level = 3
                            level = min(level, settings.pdf_max_heading_levels)
                            prefix = "#" * level + " "
                            
                        elements.append(prefix + text)
                        
        page_text_content = "\n\n".join(elements).strip()
        if page_text_content:
            pages.append(PageText(page_number=i + 1, text=page_text_content))
            
    full_text = "\n\n".join(p.text for p in pages)
    return LoadedDocument(source_file=path.name, format="pdf", text=full_text, pages=pages)
