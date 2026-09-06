"""Shared pytest fixtures.

The PDF fixture is generated on the fly with reportlab (a dev-only
dependency) rather than committed as a binary file, so the test suite has
no binary fixtures and the PDF's exact page contents always match what the
tests assert against.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from rag_api.adapters.vectorstore.embeddings import DeterministicFakeEmbeddingClient
from rag_api.adapters.vectorstore.vector_store import VectorStore

FIXTURES_DIR = Path(__file__).parent / "fixtures"


@pytest.fixture
def md_path() -> Path:
    return FIXTURES_DIR / "sample.md"


@pytest.fixture
def txt_path() -> Path:
    return FIXTURES_DIR / "sample.txt"


@pytest.fixture
def html_path() -> Path:
    return FIXTURES_DIR / "sample.html"


PDF_PAGE_1_TEXT = "Onboarding new employees begins with an equipment request submitted to IT."
PDF_PAGE_2_TEXT = "Offboarding departing employees requires revoking all system access within 24 hours."


@pytest.fixture
def pdf_path(tmp_path: Path) -> Path:
    from reportlab.pdfgen import canvas

    path = tmp_path / "sample.pdf"
    c = canvas.Canvas(str(path))
    c.drawString(72, 720, PDF_PAGE_1_TEXT)
    c.showPage()
    c.drawString(72, 720, PDF_PAGE_2_TEXT)
    c.showPage()
    c.save()
    return path


@pytest.fixture
def fake_embedder() -> DeterministicFakeEmbeddingClient:
    return DeterministicFakeEmbeddingClient(dimension=128)


@pytest.fixture
def vector_store(tmp_path: Path) -> VectorStore:
    return VectorStore(persist_dir=tmp_path / "qdrant", collection_name="test_collection", dense_dimension=128)


@pytest.fixture
def small_vector_store(tmp_path: Path) -> VectorStore:
    """3-dimensional store for tests using hand-crafted vectors like [1.0, 0.0, 0.0]."""
    return VectorStore(persist_dir=tmp_path / "qdrant_small", collection_name="test_small", dense_dimension=3)


@pytest.fixture
def docx_path(tmp_path: Path) -> Path:
    import docx
    path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_heading("Docx Heading", level=1)
    doc.add_paragraph("This is a paragraph in docx.")
    doc.add_heading("Subheading", level=2)
    doc.add_paragraph("This is a second paragraph.")
    
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Val 1"
    table.cell(1, 1).text = "Val 2"
    
    doc.save(str(path))
    return path


@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    from pptx import Presentation
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    
    # Slide 1
    slide_layout = prs.slide_layouts[0] # Title slide
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Slide 1 Title"
    slide1.shapes.placeholders[1].text = "Slide 1 Body"
    
    # Slide 2
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide 2 Title"
    slide2.shapes.placeholders[1].text = "Slide 2 Body"
    
    prs.save(str(path))
    return path


@pytest.fixture
def xlsx_path(tmp_path: Path) -> Path:
    import openpyxl
    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Col 1", "Col 2"])
    ws1.append(["Data 1", "Data 2"])
    
    ws2 = wb.create_sheet(title="Sheet2")
    ws2.append(["A", "B", "C"])
    ws2.append([1, 2, 3])
    
    wb.save(str(path))
    return path
